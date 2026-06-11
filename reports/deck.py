"""
Generate a professional presentation deck for the hackathon submission.

Creates a dark-themed PPTX with the ranking system's approach,
architecture, results, and limitations.

Usage:
    python reports/deck.py
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Colors
BG_COLOR = RGBColor(0x0F, 0x11, 0x15)
ORANGE = RGBColor(0xF7, 0x93, 0x1A)
GOLD = RGBColor(0xFF, 0xD6, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DARK_BG = RGBColor(0x03, 0x03, 0x04)


def set_slide_bg(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=13, color=WHITE, bullet_color=ORANGE):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)

    return tf


def build_deck(output_path: str = "reports/presentation.pptx"):
    """Generate the full presentation deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ── Slide 1: Title ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1.5, 2.0, 10, 1.2,
                 "Intelligent Candidate Discovery",
                 font_size=36, color=ORANGE, bold=True)
    add_text_box(slide, 1.5, 3.2, 10, 0.6,
                 "& Ranking System",
                 font_size=36, color=GOLD, bold=True)
    add_text_box(slide, 1.5, 4.5, 10, 0.5,
                 "Redrob AI Hackathon — The Data & AI Challenge",
                 font_size=18, color=MUTED)
    add_text_box(slide, 1.5, 5.3, 10, 0.4,
                 "A hybrid semantic + signal-based approach to candidate-JD matching",
                 font_size=14, color=MUTED)

    # ── Slide 2: Problem Statement ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "The Problem", font_size=28, color=ORANGE, bold=True)
    add_text_box(slide, 1.0, 1.3, 11, 0.8,
                 "Given a detailed job description and a pool of 100,000 candidate profiles, "
                 "rank the top 100 best-fit candidates with explanations.",
                 font_size=16, color=WHITE)
    add_bullet_list(slide, 1.0, 2.4, 11, 4.0, [
        "Candidates have rich structured profiles: career history, skills with proficiency levels, "
        "education, behavioral signals from the Redrob platform",
        "The JD is nuanced — it explicitly states what it wants AND what it doesn't want",
        "Simple keyword matching will fail: the dataset contains keyword-stuffer traps "
        "(non-technical roles with AI keywords listed as skills)",
        "~80 honeypot candidates with impossible profiles must be detected and excluded",
        "Behavioral signals matter: a perfect-on-paper candidate who's inactive is not actually available",
        "System must run in ≤5 minutes on CPU with 16GB RAM, no API calls",
    ], font_size=14)

    # ── Slide 3: Why Keyword Matching Fails ───────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Why Keyword Matching Fails", font_size=28, color=ORANGE, bold=True)
    add_bullet_list(slide, 1.0, 1.5, 11, 5.0, [
        "Marketing Manager with 9 AI keywords in skills list ≠ AI Engineer",
        "The JD says 'embeddings-based retrieval systems' — a candidate who writes "
        "'built a recommendation engine using dense vectors' is a match, but shares zero keywords",
        "Keyword matching can't assess career trajectory: consulting-only vs product company background",
        "It ignores behavioral signals entirely: response rate, recency, notice period",
        "It can't detect honeypots: profiles that look perfect on keywords but have impossible dates/durations",
        "The JD explicitly warns: 'The right answer is not find candidates whose skills section "
        "contains the most AI keywords. That's a trap we've built into the dataset.'",
    ], font_size=14)

    # ── Slide 4: Dataset Understanding ────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Dataset Understanding", font_size=28, color=ORANGE, bold=True)
    add_bullet_list(slide, 1.0, 1.4, 5.5, 5.5, [
        "100,000 candidate profiles in JSONL format (~487MB)",
        "Each profile contains:",
        "   ↳ Profile: headline, summary, title, company, location",
        "   ↳ Career history: 1-10 jobs with descriptions",
        "   ↳ Education: institution, degree, field, tier",
        "   ↳ Skills: name, proficiency, endorsements, duration",
        "   ↳ Certifications and languages",
        "   ↳ 23 Redrob behavioral signals",
    ], font_size=14)
    add_bullet_list(slide, 6.8, 1.4, 5.5, 5.5, [
        "Key data quality challenges:",
        "   ↳ Keyword stuffers: non-technical roles with AI keywords",
        "   ↳ ~80 honeypots with impossible profiles",
        "   ↳ Title-description mismatches (e.g., title says 'Operations' but description is about ML)",
        "   ↳ Mixed locations (India, USA, Canada, Australia)",
        "   ↳ Salary ranges in INR LPA",
        "   ↳ Skill assessments available for some candidates",
    ], font_size=14)

    # ── Slide 5: System Architecture ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "System Architecture", font_size=28, color=ORANGE, bold=True)
    add_text_box(slide, 1.0, 1.3, 11, 0.6,
                 "Two-stage pipeline: pre-computation (offline) + ranking (under 5 min)",
                 font_size=15, color=MUTED)
    add_bullet_list(slide, 1.0, 2.2, 5.5, 5.0, [
        "Stage 0 — Pre-computation (offline, one-time):",
        "   ↳ Load 100K candidates from JSONL",
        "   ↳ Build text representation per candidate",
        "   ↳ Generate 384-dim embeddings (all-MiniLM-L6-v2)",
        "   ↳ Build FAISS IndexFlatIP for exact search",
        "   ↳ Save embeddings + index to disk",
    ], font_size=14)
    add_bullet_list(slide, 6.8, 2.2, 5.5, 5.0, [
        "Stage 1 — Ranking (≤5 min, CPU, no network):",
        "   ↳ Load precomputed artifacts",
        "   ↳ Encode JD → query embedding",
        "   ↳ FAISS retrieval: top 500 by cosine similarity",
        "   ↳ Honeypot detection and filtering",
        "   ↳ Multi-signal scoring (7 dimensions)",
        "   ↳ Weighted combination → final score",
        "   ↳ Top 100, generate reasoning, export CSV",
    ], font_size=14)

    # ── Slide 6: Feature Engineering ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Scoring Dimensions", font_size=28, color=ORANGE, bold=True)
    add_bullet_list(slide, 1.0, 1.4, 5.5, 5.5, [
        "1. Semantic Similarity (25%)",
        "   ↳ Cosine sim between candidate text and JD embedding",
        "   ↳ Captures role understanding beyond keywords",
        "",
        "2. Skill Relevance (20%)",
        "   ↳ Must-have vs nice-to-have weighted matching",
        "   ↳ Trust multiplier: proficiency × duration × endorsements",
        "   ↳ Career text cross-check (do descriptions match listed skills?)",
        "",
        "3. Experience Fit (15%)",
        "   ↳ Gaussian centered on 5-9 year range",
        "   ↳ Soft penalty outside, hard penalty below 2 years",
    ], font_size=13)
    add_bullet_list(slide, 6.8, 1.4, 5.5, 5.5, [
        "4. Career Trajectory (15%)",
        "   ↳ Title relevance (AI/ML engineer vs non-technical)",
        "   ↳ Product company vs consulting-only",
        "   ↳ Job stability and production experience signals",
        "",
        "5. Behavioral Signals (15%)",
        "   ↳ Response rate, recency, notice period, GitHub",
        "   ↳ Open-to-work flag, interview completion rate",
        "",
        "6. Location (5%) + Education (5%)",
        "   ↳ India preferred, Pune/Noida bonus",
        "   ↳ CS/ML field + institution tier",
    ], font_size=13)

    # ── Slide 7: Ranking Logic Detail ─────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Ranking Logic", font_size=28, color=ORANGE, bold=True)
    add_bullet_list(slide, 1.0, 1.4, 11, 5.5, [
        "Final Score = 0.25×Semantic + 0.20×Skill + 0.15×Experience "
        "+ 0.15×Career + 0.15×Behavioral + 0.05×Location + 0.05×Education",
        "",
        "Anti-gaming measures:",
        "   ↳ Keyword stuffing detection: skills with beginner proficiency and 0 months "
        "duration get a trust penalty",
        "   ↳ Career text cross-validation: if skills list says 'embeddings' but no job "
        "description mentions anything technical, the trust score drops",
        "   ↳ Title mismatch detection: 'Marketing Manager' listing 10 AI skills → low "
        "career trajectory score",
        "",
        "Honeypot filtering:",
        "   ↳ Experience vs career history duration mismatch",
        "   ↳ Expert proficiency with zero usage months",
        "   ↳ Impossible career date ranges",
        "   ↳ 2+ flags = excluded from ranking entirely",
    ], font_size=13)

    # ── Slide 8: Explanation Layer ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Explanation Layer", font_size=28, color=ORANGE, bold=True)
    add_text_box(slide, 1.0, 1.3, 11, 0.6,
                 "Each candidate gets a specific, profile-aware reasoning — not a template.",
                 font_size=15, color=MUTED)
    add_bullet_list(slide, 1.0, 2.2, 11, 5.0, [
        'Example for a strong match:',
        '  "ML Engineer at Flipkart with 6.4 years experience; matching must-have skills: '
        'embeddings, faiss, python, ranking; strong product-company background; '
        'responsive (rate: 88%); available soon (notice: 30d); based in Bangalore."',
        '',
        'Example for a weak match:',
        '  "Marketing Manager at Wipro with 8.3 years experience; '
        'consulting-heavy background; low recruiter response rate (6%); '
        'concern: above preferred experience range; based in Chennai."',
        '',
        'Reasoning mentions: title, company, experience, matched skills, '
        'career type, response rate, notice period, location, and any concerns.',
    ], font_size=13)

    # ── Slide 9: Output Format ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Output & Validation", font_size=28, color=ORANGE, bold=True)
    add_bullet_list(slide, 1.0, 1.4, 11, 5.5, [
        "Submission format: CSV with columns candidate_id, rank, score, reasoning",
        "Exactly 100 rows, ranks 1-100, scores non-increasing",
        "Tie-breaking: candidate_id ascending",
        "",
        "Built-in validation matches the official validate_submission.py:",
        "   ↳ Row count check",
        "   ↳ Unique candidate_id and rank check",
        "   ↳ Score monotonicity check",
        "   ↳ candidate_id format check (CAND_XXXXXXX)",
        "",
        "Additional detailed output: per-dimension score breakdown CSV for analysis",
        "Demo UI: Streamlit app with JD display, candidate cards, score bars, and filters",
    ], font_size=13)

    # ── Slide 10: Limitations & Future Work ───────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1.0, 0.5, 10, 0.6,
                 "Limitations & Future Work", font_size=28, color=ORANGE, bold=True)
    add_bullet_list(slide, 1.0, 1.4, 5.5, 5.5, [
        "Current limitations:",
        "   ↳ Scoring weights are hand-tuned, not learned from data",
        "   ↳ Single embedding model — no ensemble or cross-encoder reranking",
        "   ↳ Skill matching uses string overlap, not semantic skill similarity",
        "   ↳ Honeypot detection is heuristic-based, may miss subtle cases",
        "   ↳ No multi-language support (profiles are all English)",
        "   ↳ Career description analysis uses keyword counts, not NLI",
    ], font_size=13)
    add_bullet_list(slide, 6.8, 1.4, 5.5, 5.5, [
        "What we'd improve next:",
        "   ↳ Learn-to-rank model trained on recruiter feedback data",
        "   ↳ Cross-encoder reranking stage on top-100 shortlist",
        "   ↳ Skill taxonomy with embeddings (not string matching)",
        "   ↳ Fine-tune embedding model on JD-resume pairs",
        "   ↳ Add salary range compatibility scoring",
        "   ↳ A/B test framework for weight optimization",
        "   ↳ Better honeypot detection using anomaly detection",
    ], font_size=13)

    # ── Save ──────────────────────────────────────────────────────────
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Presentation saved to {output}")


if __name__ == "__main__":
    build_deck()
